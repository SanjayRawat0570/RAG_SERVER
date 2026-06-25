"""Format-specific parsers (F9).

Each parser converts raw bytes into ``(text, metadata)``.  Lightweight formats
(text, markdown, html, json) use only the standard library.  Heavier formats
import their dependency lazily and raise ``MissingDependencyError`` when the
package is absent — keeping the rest of the system dependency-free.

Format support matrix
---------------------
Format          Library         Notes
──────────────────────────────────────────────────────────────────
.txt / .text    stdlib          Always works
.md / .markdown stdlib          Kept as-is; heading-aware chunking later
.html / .htm    stdlib          Script/style tags stripped
.json           stdlib          Flattened to key: value lines
.xml            stdlib          Text content extracted
.pdf            pypdf           Searchable text; scanned PDFs need OCR
.docx           python-docx     Paragraphs + flattened tables
.xlsx           openpyxl        Each sheet as a section
.pptx           python-pptx     Each slide as a section
.jpg/.png/.bmp  Pillow          OCR via pytesseract when installed; stub otherwise
.mp3/.wav       —               Transcription stub (whisper if installed)
.mp4/.mov       —               Transcription stub (whisper if installed)
URL (http/s)    httpx           Fetch + HTML parse
"""
from __future__ import annotations

import io
import json
from html.parser import HTMLParser
from typing import Any, Callable

ParseResult = tuple[str, dict[str, Any]]


class MissingDependencyError(RuntimeError):
    """Raised when a parser's optional dependency is not installed."""


def parse_text(content: bytes, filename: str | None) -> ParseResult:
    return content.decode("utf-8", errors="replace"), {}


def parse_markdown(content: bytes, filename: str | None) -> ParseResult:
    # Markdown is kept as-is; structure-aware chunking (F10) uses the headings.
    return content.decode("utf-8", errors="replace"), {"markdown": True}


class _HTMLTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._chunks: list[str] = []
        self._skip = 0
        self.title: str | None = None
        self._in_title = False

    _BLOCK = {"p", "div", "br", "li", "tr", "h1", "h2", "h3", "h4", "h5", "h6", "section"}
    _DROP = {"script", "style", "head", "noscript"}

    def handle_starttag(self, tag, attrs):
        if tag in self._DROP:
            self._skip += 1
        if tag == "title":
            self._in_title = True

    def handle_endtag(self, tag):
        if tag in self._DROP and self._skip:
            self._skip -= 1
        if tag == "title":
            self._in_title = False
        if tag in self._BLOCK:
            self._chunks.append("\n")

    def handle_data(self, data):
        # Capture <title> even though it lives inside the dropped <head>.
        if self._in_title:
            self.title = (self.title or "") + data.strip()
            return
        if self._skip:
            return
        if data.strip():
            self._chunks.append(data)

    def text(self) -> str:
        return "".join(self._chunks)


def parse_html(content: bytes, filename: str | None) -> ParseResult:
    parser = _HTMLTextExtractor()
    parser.feed(content.decode("utf-8", errors="replace"))
    meta: dict[str, Any] = {}
    if parser.title:
        meta["title"] = parser.title
    return parser.text(), meta


def parse_json(content: bytes, filename: str | None) -> ParseResult:
    data = json.loads(content.decode("utf-8", errors="replace"))
    lines: list[str] = []

    def walk(node: Any, prefix: str = "") -> None:
        if isinstance(node, dict):
            for k, v in node.items():
                walk(v, f"{prefix}{k}: " if not isinstance(v, (dict, list)) else f"{prefix}{k}.")
        elif isinstance(node, list):
            for item in node:
                walk(item, prefix)
        else:
            lines.append(f"{prefix}{node}")

    walk(data)
    return "\n".join(lines), {"structured": True}


def parse_pdf(content: bytes, filename: str | None) -> ParseResult:
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover - dependency always present here
        raise MissingDependencyError("PDF parsing requires 'pypdf'") from exc
    import io

    reader = PdfReader(io.BytesIO(content))
    pages = [page.extract_text() or "" for page in reader.pages]
    meta: dict[str, Any] = {"pages": len(pages)}
    if reader.metadata and reader.metadata.title:
        meta["title"] = reader.metadata.title
    return "\n\n".join(pages), meta


def parse_docx(content: bytes, filename: str | None) -> ParseResult:
    try:
        import docx
    except ImportError as exc:
        raise MissingDependencyError("DOCX parsing requires 'python-docx'") from exc

    document = docx.Document(io.BytesIO(content))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    # Flatten tables to tab-separated rows so their content is retained.
    for table in document.tables:
        for row in table.rows:
            paragraphs.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(paragraphs), {}


def parse_xlsx(content: bytes, filename: str | None) -> ParseResult:
    """Excel spreadsheet — each sheet becomes a '### Sheet: name' section."""
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise MissingDependencyError("XLSX parsing requires 'openpyxl'") from exc

    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    sections: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            row_text = "\t".join("" if cell is None else str(cell) for cell in row)
            if row_text.strip():
                rows.append(row_text)
        if rows:
            sections.append(f"### Sheet: {sheet_name}\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sections), {"sheets": list(wb.sheetnames)}


def parse_pptx(content: bytes, filename: str | None) -> ParseResult:
    """PowerPoint — each slide becomes a '### Slide N' section."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise MissingDependencyError("PPTX parsing requires 'python-pptx'") from exc

    prs = Presentation(io.BytesIO(content))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"### Slide {i}\n" + "\n".join(texts))
    return "\n\n".join(slides), {"slides": len(prs.slides)}


def parse_xml(content: bytes, filename: str | None) -> ParseResult:
    """XML — extract all text nodes, stripping tags."""
    import xml.etree.ElementTree as ET

    try:
        root = ET.fromstring(content)
    except ET.ParseError as exc:
        raise ValueError(f"Malformed XML: {exc}") from exc

    parts: list[str] = []

    def _walk(node: ET.Element) -> None:
        if node.text and node.text.strip():
            parts.append(node.text.strip())
        for child in node:
            _walk(child)
        if node.tail and node.tail.strip():
            parts.append(node.tail.strip())

    _walk(root)
    return "\n".join(parts), {"root_tag": root.tag}


def parse_image(content: bytes, filename: str | None) -> ParseResult:
    """Image file — OCR via pytesseract when installed; informative stub otherwise.

    Uses Pillow to read basic image properties (always), then tries pytesseract
    for OCR text extraction.  When pytesseract / Tesseract is unavailable the
    result is a structured stub so downstream quality assessment can flag it.
    """
    meta: dict[str, Any] = {"ocr_attempted": False}

    # Basic image metadata (Pillow is always available).
    try:
        from PIL import Image as _PILImage
        img = _PILImage.open(io.BytesIO(content))
        meta.update({"width": img.width, "height": img.height, "mode": img.mode})
    except Exception:
        pass

    # Attempt OCR.
    try:
        import pytesseract
        from PIL import Image as _PILImage
        img = _PILImage.open(io.BytesIO(content))
        text = pytesseract.image_to_string(img)
        meta["ocr_attempted"] = True
        meta["ocr_engine"] = "pytesseract"
        return text, meta
    except ImportError:
        pass
    except Exception as exc:
        meta["ocr_error"] = str(exc)

    # Stub: OCR unavailable.
    size = f"{meta.get('width', '?')}×{meta.get('height', '?')}"
    stub = (
        f"[Image file: {filename or 'unknown'}, {size} pixels. "
        "OCR text extraction requires pytesseract + Tesseract to be installed. "
        "Install with: pip install pytesseract && apt-get install tesseract-ocr]"
    )
    meta["ocr_available"] = False
    meta["quality_flag"] = "ocr_unavailable"
    return stub, meta


def parse_audio(content: bytes, filename: str | None) -> ParseResult:
    """Audio file — transcription via OpenAI Whisper when installed; stub otherwise.

    Whisper runs fully offline on CPU (no API key needed). Install with:
        pip install openai-whisper
    """
    import tempfile, os

    meta: dict[str, Any] = {"transcription_attempted": False, "file_size_bytes": len(content)}
    ext = (filename or "audio.mp3").rsplit(".", 1)[-1].lower()

    try:
        import whisper  # type: ignore[import]

        with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
            tmp.write(content)
            tmp_path = tmp.name
        try:
            model = whisper.load_model("tiny")
            result = model.transcribe(tmp_path)
            text = result.get("text", "")
            meta["transcription_attempted"] = True
            meta["transcription_engine"] = "whisper"
            meta["language"] = result.get("language")
        finally:
            os.unlink(tmp_path)
        return text, meta

    except ImportError:
        pass
    except Exception as exc:
        meta["transcription_error"] = str(exc)

    stub = (
        f"[Audio file: {filename or 'audio'} ({len(content):,} bytes). "
        "Transcription requires OpenAI Whisper. "
        "Install with: pip install openai-whisper]"
    )
    meta["transcription_available"] = False
    meta["quality_flag"] = "transcription_unavailable"
    return stub, meta


async def fetch_url(url: str) -> ParseResult:
    """Fetch a web page and extract clean text (uses httpx, always available)."""
    try:
        import httpx as _httpx
    except ImportError as exc:
        raise MissingDependencyError("URL fetching requires 'httpx'") from exc

    async with _httpx.AsyncClient(
        timeout=20,
        follow_redirects=True,
        headers={"User-Agent": "RAG-Ingest-Bot/1.0 (document indexing)"},
    ) as client:
        resp = await client.get(url)

    resp.raise_for_status()
    parser = _HTMLTextExtractor()
    parser.feed(resp.text)
    meta: dict[str, Any] = {
        "url":         url,
        "status_code": resp.status_code,
        "content_type": resp.headers.get("content-type", ""),
    }
    if parser.title:
        meta["title"] = parser.title
    return parser.text(), meta


Parser = Callable[[bytes, "str | None"], ParseResult]
