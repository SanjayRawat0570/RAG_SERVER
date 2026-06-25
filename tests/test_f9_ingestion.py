"""Tests for F9: Document Ingestion & Format Support."""
from __future__ import annotations

import io
import json
import textwrap

import pytest


# ── Helpers ────────────────────────────────────────────────────────────────────

def _xlsx_bytes() -> bytes:
    """Create a minimal in-memory XLSX file."""
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Revenue", "Year"])
    ws.append(["ACME",  1_000_000, 2024])
    ws.append(["Corp",  2_500_000, 2024])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _pptx_bytes() -> bytes:
    """Create a minimal in-memory PPTX file."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Revenue Report"
    slide.placeholders[1].text = "Q3 revenue increased by 20 percent."
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _png_bytes() -> bytes:
    """Create a 10×10 white PNG image."""
    from PIL import Image
    img = Image.new("RGB", (10, 10), color=(255, 255, 255))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _pdf_bytes() -> bytes:
    """Create a minimal PDF-like bytes (won't parse but tests MissingDep path)."""
    # Use the existing pypdf to generate real PDF bytes
    try:
        from pypdf import PdfWriter
        writer = PdfWriter()
        writer.add_blank_page(width=200, height=200)
        buf = io.BytesIO()
        writer.write(buf)
        return buf.getvalue()
    except Exception:
        return b"%PDF-1.4 minimal"


# ── Parser unit tests ─────────────────────────────────────────────────────────

def test_f9_parse_text():
    from app.rag.ingestion.parsers import parse_text
    text, meta = parse_text(b"Hello world!", "doc.txt")
    assert "Hello" in text
    assert meta == {}


def test_f9_parse_markdown():
    from app.rag.ingestion.parsers import parse_markdown
    text, meta = parse_markdown(b"# Title\n\nBody text.", "doc.md")
    assert "Title" in text
    assert meta.get("markdown") is True


def test_f9_parse_html():
    from app.rag.ingestion.parsers import parse_html
    html = b"<html><head><title>Test</title></head><body><p>Hello</p></body></html>"
    text, meta = parse_html(html, "doc.html")
    assert "Hello" in text
    assert meta.get("title") == "Test"


def test_f9_parse_html_strips_script():
    from app.rag.ingestion.parsers import parse_html
    html = b"<html><body><script>var x=1;</script><p>Clean</p></body></html>"
    text, _ = parse_html(html, "doc.html")
    assert "var x=1" not in text
    assert "Clean" in text


def test_f9_parse_json():
    from app.rag.ingestion.parsers import parse_json
    data = {"company": "ACME", "revenue": 1000000}
    text, meta = parse_json(json.dumps(data).encode(), "data.json")
    assert "ACME" in text
    assert "revenue" in text
    assert meta.get("structured") is True


def test_f9_parse_xml():
    from app.rag.ingestion.parsers import parse_xml
    xml = b"<root><company>ACME</company><value>100</value></root>"
    text, meta = parse_xml(xml, "data.xml")
    assert "ACME" in text
    assert "100" in text
    assert "root" in meta.get("root_tag", "root")


def test_f9_parse_docx():
    from app.rag.ingestion.parsers import parse_docx
    import docx, io as _io
    doc = docx.Document()
    doc.add_paragraph("Quarterly results were positive.")
    buf = _io.BytesIO()
    doc.save(buf)
    text, _ = parse_docx(buf.getvalue(), "doc.docx")
    assert "Quarterly" in text


def test_f9_parse_xlsx():
    from app.rag.ingestion.parsers import parse_xlsx
    text, meta = parse_xlsx(_xlsx_bytes(), "data.xlsx")
    assert "ACME" in text
    assert "Revenue" in text
    assert "Sheet1" in meta.get("sheets", [])


def test_f9_parse_xlsx_multiple_sheets():
    from openpyxl import Workbook
    from app.rag.ingestion.parsers import parse_xlsx
    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Sales"
    ws1.append(["Q1", 100])
    ws2 = wb.create_sheet("Costs")
    ws2.append(["Q1", 50])
    buf = io.BytesIO()
    wb.save(buf)
    text, meta = parse_xlsx(buf.getvalue(), "data.xlsx")
    assert "### Sheet: Sales" in text
    assert "### Sheet: Costs" in text
    assert "Sales" in meta["sheets"]


def test_f9_parse_pptx():
    from app.rag.ingestion.parsers import parse_pptx
    text, meta = parse_pptx(_pptx_bytes(), "slides.pptx")
    assert "Revenue" in text
    assert isinstance(meta.get("slides"), int)
    assert meta["slides"] >= 1


def test_f9_parse_pptx_slide_sections():
    from app.rag.ingestion.parsers import parse_pptx
    text, _ = parse_pptx(_pptx_bytes(), "slides.pptx")
    assert "### Slide 1" in text


def test_f9_parse_image_stub_when_no_ocr():
    from app.rag.ingestion.parsers import parse_image
    text, meta = parse_image(_png_bytes(), "screenshot.png")
    # pytesseract not installed → stub returned
    assert isinstance(text, str) and len(text) > 0
    if not meta.get("ocr_attempted"):
        assert "ocr_available" in meta or "OCR" in text


def test_f9_parse_image_metadata():
    from app.rag.ingestion.parsers import parse_image
    _, meta = parse_image(_png_bytes(), "screenshot.png")
    # Pillow should always provide dimensions
    assert meta.get("width") == 10
    assert meta.get("height") == 10


def test_f9_parse_audio_stub():
    from app.rag.ingestion.parsers import parse_audio
    text, meta = parse_audio(b"\x00\x01\x02", "audio.mp3")
    # Whisper not installed → stub
    assert isinstance(text, str)
    if not meta.get("transcription_attempted"):
        assert "transcription_available" in meta or "Transcription" in text


# ── Registry tests ────────────────────────────────────────────────────────────

def test_f9_detect_format_from_extension():
    from app.rag.ingestion.registry import detect_format
    assert detect_format("report.pdf") == "pdf"
    assert detect_format("data.xlsx") == "xlsx"
    assert detect_format("slides.pptx") == "pptx"
    assert detect_format("photo.jpg") == "jpg"
    assert detect_format("photo.jpeg") == "jpeg"
    assert detect_format("audio.mp3") == "mp3"
    assert detect_format("video.mp4") == "mp4"


def test_f9_detect_format_explicit_override():
    from app.rag.ingestion.registry import detect_format
    assert detect_format("noext", "pdf") == "pdf"


def test_f9_detect_format_unknown_falls_back():
    from app.rag.ingestion.registry import detect_format
    assert detect_format("file.unknown") == "text"


def test_f9_supported_formats_covers_all_new():
    from app.rag.ingestion.registry import supported_formats
    fmts = supported_formats()
    assert "xlsx"  in fmts
    assert "pptx"  in fmts
    assert "image" in fmts
    assert "audio" in fmts or "video" in fmts
    assert "xml"   in fmts


def test_f9_ingest_xlsx():
    from app.rag.ingestion.registry import ingest
    doc = ingest(_xlsx_bytes(), filename="data.xlsx")
    assert "ACME" in doc.text
    assert doc.format == "xlsx"
    assert "quality" in doc.metadata


def test_f9_ingest_pptx():
    from app.rag.ingestion.registry import ingest
    doc = ingest(_pptx_bytes(), filename="slides.pptx")
    assert "Revenue" in doc.text
    assert doc.format == "pptx"
    assert "quality" in doc.metadata


def test_f9_ingest_attaches_quality():
    from app.rag.ingestion.registry import ingest
    doc = ingest(b"The quick brown fox jumps over the lazy dog. " * 5, filename="doc.txt")
    q = doc.metadata["quality"]
    assert "score" in q
    assert "action" in q
    assert "language" in q


# ── Quality assessment tests ──────────────────────────────────────────────────

def test_f9_quality_good_text():
    from app.rag.ingestion.quality import assess_quality
    text = "The quarterly revenue report shows strong growth across all divisions. " * 5
    r = assess_quality(text)
    assert r.score >= 0.7
    assert r.action == "ok"
    assert r.readable is True


def test_f9_quality_too_short():
    from app.rag.ingestion.quality import assess_quality
    r = assess_quality("hi")
    assert "too_short" in r.flags
    assert r.score < 0.7


def test_f9_quality_very_short():
    from app.rag.ingestion.quality import assess_quality
    r = assess_quality("hello world this is text")
    assert "very_short" in r.flags or "too_short" in r.flags


def test_f9_quality_ocr_noise():
    from app.rag.ingestion.quality import assess_quality
    # high proportion of weird chars
    noise = "##$$%%^^&&**((||\\//~~``@@!!" * 20
    r = assess_quality(noise)
    assert "ocr_noise" in r.flags or r.score < 0.7


def test_f9_quality_score_range():
    from app.rag.ingestion.quality import assess_quality
    r = assess_quality("anything goes here" * 3)
    assert 0.0 <= r.score <= 1.0


def test_f9_quality_action_warn_triggers_note():
    from app.rag.ingestion.quality import assess_quality
    short = "word " * 8  # very_short
    r = assess_quality(short)
    if r.action == "warn":
        assert r.note is not None


def test_f9_quality_language_detected():
    from app.rag.ingestion.quality import assess_quality
    text = "The quarterly revenue report shows strong growth. " * 10
    r = assess_quality(text)
    if r.language:
        assert len(r.language) <= 5  # ISO code, e.g. "en"
        assert 0.0 <= r.language_confidence <= 1.0


def test_f9_quality_ocr_unavailable_flag():
    from app.rag.ingestion.quality import assess_quality
    stub = "[Image file: photo.png, 100×100. OCR unavailable.]"
    r = assess_quality(stub, {"quality_flag": "ocr_unavailable"})
    assert "ocr_unavailable" in r.flags


# ── Cleaning tests ────────────────────────────────────────────────────────────

def test_f9_clean_text_normalizes_whitespace():
    from app.rag.ingestion.cleaning import clean_text
    text = "Hello   world\r\n\r\nFoo  bar"
    cleaned = clean_text(text)
    assert "   " not in cleaned
    assert "\r" not in cleaned


def test_f9_clean_text_collapses_blank_lines():
    from app.rag.ingestion.cleaning import clean_text
    text = "Para 1\n\n\n\n\nPara 2"
    cleaned = clean_text(text)
    assert "\n\n\n" not in cleaned


def test_f9_derive_title_from_heading():
    from app.rag.ingestion.cleaning import derive_title
    text = "# Annual Report 2024\n\nBody text."
    assert derive_title(text) == "Annual Report 2024"


def test_f9_derive_title_first_line_fallback():
    from app.rag.ingestion.cleaning import derive_title
    text = "Summary of quarterly results\n\nBody text."
    assert derive_title(text) == "Summary of quarterly results"


# ── API endpoint tests ────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_f9_api_formats():
    """GET /ingest/formats returns all expected format keys."""
    from fastapi.testclient import TestClient
    from app.main import app

    with TestClient(app) as client:
        resp = client.get(
            "/api/v1/ingest/formats",
            headers={"Authorization": "Bearer dev"},
        )
    assert resp.status_code == 200
    data = resp.json()
    fmts = data["formats"]
    assert "xlsx"  in fmts
    assert "pptx"  in fmts
    assert "image" in fmts
    assert "pdf"   in fmts
    assert "docx"  in fmts


@pytest.mark.asyncio
async def test_f9_api_analyze_text_file():
    """POST /ingest/analyze with a .txt file returns quality + preview."""
    from fastapi.testclient import TestClient
    from app.main import app

    content = b"The quarterly revenue report shows strong growth across all divisions. " * 5

    with TestClient(app) as client:
        resp = client.post(
            "/api/v1/ingest/analyze",
            files={"file": ("report.txt", content, "text/plain")},
            headers={"Authorization": "Bearer dev"},
        )
    assert resp.status_code == 200
    data = resp.json()
    assert data["format"] == "text"
    assert data["word_count"] > 0
    assert "quality" in data
    assert "preview" in data


@pytest.mark.asyncio
async def test_f9_api_analyze_xlsx():
    """POST /ingest/analyze with an .xlsx file parses sheets."""
    from fastapi.testclient import TestClient
    from app.main import app

    with TestClient(app) as client:
        resp = client.post(
            "/api/v1/ingest/analyze",
            files={"file": ("data.xlsx", _xlsx_bytes(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")},
            headers={"Authorization": "Bearer dev"},
        )
    assert resp.status_code == 200
    data = resp.json()
    assert data["format"] == "xlsx"
    assert "ACME" in data["preview"]


@pytest.mark.asyncio
async def test_f9_api_analyze_pptx():
    """POST /ingest/analyze with a .pptx file extracts slide text."""
    from fastapi.testclient import TestClient
    from app.main import app

    with TestClient(app) as client:
        resp = client.post(
            "/api/v1/ingest/analyze",
            files={"file": ("slides.pptx", _pptx_bytes(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
            headers={"Authorization": "Bearer dev"},
        )
    assert resp.status_code == 200
    data = resp.json()
    assert data["format"] == "pptx"
    assert data["word_count"] > 0


@pytest.mark.asyncio
async def test_f9_api_analyze_image():
    """POST /ingest/analyze with a .png file returns image metadata."""
    from fastapi.testclient import TestClient
    from app.main import app

    with TestClient(app) as client:
        resp = client.post(
            "/api/v1/ingest/analyze",
            files={"file": ("screenshot.png", _png_bytes(), "image/png")},
            headers={"Authorization": "Bearer dev"},
        )
    assert resp.status_code == 200
    data = resp.json()
    assert data["format"] == "image"
    assert "quality" in data


@pytest.mark.asyncio
async def test_f9_api_ingest_text():
    """POST /ingest/text indexes raw text and returns quality info."""
    from fastapi.testclient import TestClient
    from app.main import app

    body = {
        "text": "The quarterly revenue shows growth across all business units. " * 10,
        "filename": "report.txt",
        "tenant": "f9test",
        "index": True,
    }
    with TestClient(app) as client:
        resp = client.post(
            "/api/v1/ingest/text",
            json=body,
            headers={"Authorization": "Bearer dev"},
        )
    assert resp.status_code == 200
    data = resp.json()
    assert data["word_count"] > 0
    assert "quality" in data
    assert data.get("indexed") is True


@pytest.mark.asyncio
async def test_f9_api_ingest_text_no_index():
    """POST /ingest/text with index=False skips indexing."""
    from fastapi.testclient import TestClient
    from app.main import app

    with TestClient(app) as client:
        resp = client.post(
            "/api/v1/ingest/text",
            json={"text": "Hello world this is a test.", "index": False},
            headers={"Authorization": "Bearer dev"},
        )
    assert resp.status_code == 200
    data = resp.json()
    assert data.get("indexed") is False


@pytest.mark.asyncio
async def test_f9_api_upload_xlsx():
    """POST /documents/upload now accepts .xlsx files."""
    from fastapi.testclient import TestClient
    from app.main import app

    with TestClient(app) as client:
        resp = client.post(
            "/api/v1/documents/upload",
            files={"file": ("data.xlsx", _xlsx_bytes(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")},
            data={"tenant": "f9upload"},
            headers={"Authorization": "Bearer dev"},
        )
    assert resp.status_code == 200
    data = resp.json()
    assert "execution_id" in data
